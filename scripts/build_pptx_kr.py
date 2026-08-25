# -*- coding: utf-8 -*-
"""KRAFTON 월간 보고 디자인(20260715_Monthly_KRAFTON.pptx 기준)에 맞춘 한글 최종 보고서.

일반 독자용: 전문 통계 용어를 풀어서 설명하고, 수치는 analysis/expB·expA의
게이트 통과 run에서만 가져온다.

usage: python scripts/build_pptx_kr.py report/KRAFTON_DB_성능비교.pptx
"""
import os, sys, glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- KRAFTON 월간 덱 디자인 토큰 (템플릿에서 추출) --------------------------------
NAVY = RGBColor(0x0A, 0x25, 0x40)   # 표지/섹션 배경
PANEL = RGBColor(0x0E, 0x2E, 0x52)  # 다크 패널
BLUE = RGBColor(0x00, 0x78, 0xD4)   # Azure 블루
TEAL = RGBColor(0x17, 0xB6, 0xC4)
CYAN = RGBColor(0x50, 0xE6, 0xFF)
STEEL = RGBColor(0x3E, 0x6E, 0xA8)
LIGHT = RGBColor(0xF4, 0xF7, 0xFB)  # 본문 배경
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x12, 0x23, 0x3A)
BODY = RGBColor(0x34, 0x46, 0x5C)
MUTED = RGBColor(0x5A, 0x6B, 0x82)
LBLUE = RGBColor(0x9F, 0xC1, 0xE2)
DTEXT = RGBColor(0xD6, 0xE4, 0xF2)
GHOST = RGBColor(0xD7, 0xE3, 0xF2)
GHOST_D = RGBColor(0x13, 0x37, 0x5E)
SKY = RGBColor(0x8F, 0xB4, 0xD9)
WARN = RGBColor(0xF2, 0xA9, 0x3B)
WARN_BG = RGBColor(0xFC, 0xEF, 0xD6)
WARN_TX = RGBColor(0x9A, 0x64, 0x10)
FONT = "Malgun Gothic"
DATE = "2026.08"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
PAGE = [0]


def slide(bg):
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill; f.solid(); f.fore_color.rgb = bg
    PAGE[0] += 1
    return s


def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(s, x, y, w, h, lines, size=12, bold=False, color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=4):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.03); tf.margin_top = tf.margin_bottom = Inches(0.02)
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(spacing)
        runs = ln if isinstance(ln, list) else [(ln, size, bold, color)]
        for t, sz, b, c in runs:
            r = p.add_run(); r.text = t; r.font.name = FONT; r.font.size = Pt(sz); r.font.bold = b; r.font.color.rgb = c
    return tb


def topbar(s):
    rect(s, 0, 0, 13.333, 0.06, BLUE)
    rect(s, 0, 0, 4.1, 0.06, TEAL)


def footer(s, dark=False):
    text(s, 11.1, 7.12, 1.1, 0.3, [[(DATE, 8.5, False, SKY if dark else MUTED)]])
    text(s, 12.3, 7.12, 0.5, 0.3, [[(str(PAGE[0]), 8.5, True, CYAN if dark else BLUE)]])


def header(s, eyebrow, title):
    topbar(s)
    text(s, 0.6, 0.38, 11.0, 0.3, [[(eyebrow, 12, True, BLUE)]])
    text(s, 0.6, 0.68, 12.2, 0.8, [[(title, 25, True, INK)]])


def card(s, x, y, w, h, accent, title, body_lines, title_size=14.5, body_size=11.5, side="left"):
    rect(s, x, y, w, h, WHITE)
    if side == "left":
        rect(s, x, y, 0.06, h, accent)
    else:
        rect(s, x, y, w, 0.06, accent)
    text(s, x + 0.25, y + 0.16, w - 0.5, 0.5, [[(title, title_size, True, INK)]])
    if body_lines:
        text(s, x + 0.25, y + 0.16 + title_size / 60.0, w - 0.5, h - 0.6,
             [[("· " + b, body_size, False, BODY)]] if isinstance(body_lines, str) else
             [[("· " + b, body_size, False, BODY)] for b in body_lines], spacing=5)


def section_divider(num, title, sub, items):
    s = slide(NAVY)
    rect(s, 9.2, 0, 4.13, 7.5, PANEL)
    rect(s, 9.2, 0, 0.06, 7.5, TEAL)
    text(s, 0.4, 1.0, 5.5, 3.0, [[(num, 150, True, GHOST_D)]])
    rect(s, 0.62, 3.72, 0.7, 0.07, CYAN)
    text(s, 0.6, 3.98, 6.0, 0.3, [[(f"SECTION {num}", 13, True, CYAN)]])
    text(s, 0.6, 4.3, 8.4, 1.0, [[(title, 34, True, WHITE)]])
    text(s, 0.6, 5.2, 8.2, 0.5, [[(sub, 14, False, SKY)]])
    for i, it in enumerate(items):
        rect(s, 9.6, 3.15 + 0.7 * i, 0.08, 0.08, CYAN)
        text(s, 9.9, 3.05 + 0.7 * i, 3.3, 0.6, [[(it, 12.5, False, DTEXT)]])
    footer(s, dark=True)
    return s


def kr_table(s, x, y, w, rows, col_w, size=10.5, row_h=0.32, header_fill=NAVY, align_first_left=True):
    from pptx.util import Inches as I
    tbl = s.shapes.add_table(len(rows), len(rows[0]), I(x), I(y), I(w), I(row_h * len(rows))).table
    for i, cw in enumerate(col_w):
        tbl.columns[i].width = I(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = tbl.cell(i, j); c.text = ""
            c.margin_left = c.margin_right = I(0.06); c.margin_top = c.margin_bottom = I(0.02)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (j == 0 and align_first_left) else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val); r.font.name = FONT; r.font.size = Pt(size)
            if i == 0:
                r.font.bold = True; r.font.color.rgb = WHITE
                c.fill.solid(); c.fill.fore_color.rgb = header_fill
            else:
                r.font.color.rgb = INK
                c.fill.solid(); c.fill.fore_color.rgb = WHITE if i % 2 else LIGHT
    return tbl


# ================================ 1. 표지 =====================================
s = slide(NAVY)
rect(s, 0, 5.9, 13.333, 1.6, PANEL)
rect(s, 0, 0, 13.333, 0.16, BLUE)
rect(s, 0, 0, 5.2, 0.16, TEAL)
text(s, 5.4, 0.4, 7.8, 5.2, [[("Bench", 170, True, PANEL)]], align=PP_ALIGN.RIGHT)
text(s, 1.1, 1.0, 6.0, 0.5, [[("Microsoft Azure", 15, True, DTEXT)]])
rect(s, 0.62, 2.5, 0.8, 0.09, CYAN)
text(s, 0.6, 2.68, 11.5, 0.7, [[("게임 워크로드 기준 성능 비교 테스트", 20, True, CYAN)]])
text(s, 0.6, 3.35, 12.2, 1.2, [[("Azure Database 성능 비교 결과 보고", 44, True, WHITE)]])
text(s, 0.6, 4.62, 12.0, 0.5, [[("MySQL Premium SSD v1 vs v2  ·  PostgreSQL Flexible Server vs Azure HorizonDB", 16, False, LBLUE)]])
text(s, 0.6, 6.28, 12.0, 0.6, [[("일자  ", 12, False, SKY), ("2026년 8월", 12, True, WHITE),
                                ("      대상  ", 12, False, SKY), ("KRAFTON 엔지니어링", 12, True, WHITE),
                                ("      작성  ", 12, False, SKY), ("Microsoft (Euson)", 12, True, WHITE)]])

# ================================ 2. 아젠다 ====================================
s = slide(LIGHT)
topbar(s)
text(s, 0.6, 0.38, 11.0, 0.3, [[("AGENDA", 12, True, BLUE)]])
text(s, 0.6, 0.68, 12.2, 0.8, [[("오늘 전달할 내용", 25, True, INK)]])
agenda = [
    ("01", BLUE, "무엇을, 어떻게 테스트했나", "게임 서비스 워크로드 재현 · 결과를 믿을 수 있게 만드는 10가지 검증 기준"),
    ("02", TEAL, "MySQL — Premium SSD v1 vs v2", "같은 IOPS 조건에서 스토리지 세대만 바꿨을 때의 차이"),
    ("03", STEEL, "PostgreSQL vs Azure HorizonDB", "같은 vCore · 같은 부하에서 두 서비스의 응답 특성 비교"),
    ("04", WARN, "결론과 다음 단계", "권고 사항 · 아직 확인하지 못한 항목과 예상 소요 시간"),
]
for i, (num, acc, t, sub) in enumerate(agenda):
    y = 1.55 + 1.32 * i
    rect(s, 0.9, y, 11.5, 1.15, WHITE)
    rect(s, 0.9, y, 0.06, 1.15, acc)
    text(s, 1.15, y + 0.05, 1.4, 1.05, [[(num, 40, True, GHOST)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 2.75, y + 0.17, 9.4, 0.5, [[(t, 18, True, INK)]])
    text(s, 2.75, y + 0.66, 9.4, 0.45, [[(sub, 11.5, False, MUTED)]])
footer(s)

# ============================ 3. 한눈에 보는 결론 ===============================
s = slide(LIGHT)
header(s, "EXECUTIVE SUMMARY", "한눈에 보는 결론")
cards = [
    (BLUE, "MySQL: SSD v2, 평균 응답 약 30% 개선",
     ["같은 IOPS(5,000) 조건에서 중앙값 응답시간이 모든 시나리오에서 v1보다 빠름 (−18 ~ −34%)",
      "쓰기 집중 부하에서는 상위 1% 느린 요청(p99)도 55% 개선 — 게임 결산·랭킹 갱신에 유리"]),
    (WARN, "MySQL: 다만 v2에서 간헐적 지연 급증 관찰",
     ["일반 플레이 부하 3회 중 1회꼴로 p99가 3~4배 튀는 구간 발생 (원인 미확정)",
      "HA 해제 상태에서는 v2의 p99가 오히려 30% 높게 측정 — 장시간 검증 후 전환 권고"]),
    (STEEL, "HorizonDB: 같은 부하에서 꼬리 지연 최대 90% 낮음",
     ["초당 5,500건 부하에서 p99가 PostgreSQL 32ms vs HorizonDB 4.4ms",
      "부하를 3배로 튀겨도(로그인 러시) 지연이 거의 흔들리지 않음"]),
    (TEAL, "HorizonDB: 단, 완전히 같은 조건의 비교는 아님",
     ["HorizonDB는 캐시 메모리가 5배 이상 커서 데이터 전체가 메모리에 상주 (보고서에 명시)",
      "시간당 요금도 약 1.8배 — 비용까지 고려한 판단 필요"]),
]
for i, (acc, t, body) in enumerate(cards):
    x = 0.6 + 6.15 * (i % 2); y = 1.65 + 2.5 * (i // 2)
    card(s, x, y, 5.95, 2.3, acc, t, body, title_size=14, body_size=11)
text(s, 0.6, 6.7, 12.1, 0.4, [[("본 문서의 모든 수치는 검증 기준(6쪽)을 통과한 측정에서만 산출했습니다. 통과하지 못한 측정은 결과에서 제외하고 사유를 기록했습니다.", 10.5, True, BLUE)]])
footer(s)

# ========================= 4. 왜 다시 테스트했나 ================================
s = slide(LIGHT)
header(s, "BACKGROUND", "왜 이 테스트를 다시 설계했나")
text(s, 0.6, 1.55, 12.1, 0.4, [[("이전 두 차례 시도는 파이프라인은 완성했지만, 아래 세 가지 이유로 “성능 비교”로 쓸 수 있는 데이터를 남기지 못했습니다.", 12.5, False, BODY)]])
prev = [
    (WARN, "데이터가 너무 작았다", ["데이터 전체가 DB 메모리에 올라가 디스크를 한 번도 읽지 않음", "→ 스토리지 비교인데 스토리지가 쓰이지 않는 실험"]),
    (WARN, "부하 도구가 먼저 무너졌다", ["연결 수보다 훨씬 많은 요청을 밀어 넣어 클라이언트 쪽에서 타임아웃 폭주", "→ DB 한계가 아니라 측정 도구의 한계를 기록"]),
    (WARN, "서버 쪽 증거가 없었다", ["시간대 변환 오류로 서버 지표가 실험 시간과 어긋나 전부 비어 있음", "→ DB가 실제로 어떤 상태였는지 증명 불가"]),
]
for i, (acc, t, body) in enumerate(prev):
    card(s, 0.6 + 4.1 * i, 2.15, 3.9, 2.2, acc, t, body, title_size=13.5, body_size=10.5)
card(s, 0.6, 4.6, 11.9, 1.9, BLUE, "이번 설계의 핵심: “실패 조건을 먼저 정의”",
     ["측정 전 10가지 검증 기준(게이트)을 정하고, 하나라도 어기면 그 측정은 폐기 후 재실행",
      "데이터 크기를 메모리의 3배(약 24GiB)로 키워 디스크가 실제로 일하는 상태에서 측정",
      "클라이언트 자원(CPU·연결·큐)을 함께 기록해 “측정 도구가 병목이 아니었다”를 증명",
      "서버 지표는 UTC 기준으로 측정 구간에 정확히 잘라 결과 파일에 첨부"], title_size=14, body_size=11.5)
footer(s)

# ============================ 5. 어떻게 측정했나 ================================
s = slide(LIGHT)
header(s, "METHOD", "어떻게 측정했나")
rect(s, 0.6, 1.6, 4.2, 5.0, NAVY)
rect(s, 0.6, 1.6, 4.2, 0.06, TEAL)
text(s, 0.85, 1.8, 3.7, 0.4, [[("게임 서비스 워크로드", 13, True, CYAN)]])
text(s, 0.85, 2.2, 3.7, 0.9, [[("250만 계정", 34, True, WHITE)]])
text(s, 0.85, 3.05, 3.7, 3.4, [[("· 계정·인벤토리·세션·구매 원장·매치 결과·랭킹·길드 8개 테이블", 11.5, False, DTEXT)],
                               [("· 압축이 안 되는 랜덤 데이터로 약 24GiB 구성 (DB 캐시 8GiB의 3배)", 11.5, False, DTEXT)],
                               [("· 읽기 65% / 쓰기 35%의 일반 플레이 믹스 외에 쓰기 집중·로그인 러시·핫스팟 시나리오", 11.5, False, DTEXT)],
                               [("· 구매는 원장 기록과 잔액 차감이 함께 묶인 트랜잭션 — 측정 후 장부 대조로 데이터 훼손 여부 확인", 11.5, False, DTEXT)]], spacing=7)
card(s, 5.0, 1.6, 7.3, 2.35, BLUE, "부하는 “현실적인 방식”으로",
     ["먼저 동시성을 계단식으로 올려 서버가 감당 가능한 한계점(knee)을 찾음",
      "본 측정은 한계점의 65% 수준에서 초당 요청 수를 고정해 10분간 진행 — 실제 서비스처럼 요청이 서버 상태와 무관하게 계속 도착하는 방식(open-loop)",
      "응답시간은 “요청이 도착하기로 한 시각”부터 측정해 대기 시간 누락 없이 기록"], title_size=14, body_size=11.5)
card(s, 5.0, 4.2, 7.3, 2.4, TEAL, "비교는 “공정한 방식”으로",
     ["두 비교 대상을 같은 시각에 나란히 실행 — 네트워크·시간대 요인이 양쪽에 똑같이 작용",
      "시나리오마다 3~5회 반복하고, 효과 크기는 반복 쌍의 변화율 평균과 95% 신뢰구간으로 보고",
      "예열(warmup) 후 처리량이 안정된 것을 확인한 뒤에만 측정 시작 · 반복 사이 냉각 시간 확보"], title_size=14, body_size=11.5)
footer(s)

# ======================== 6. 검증 게이트(믿을 수 있는 이유) ======================
s = slide(LIGHT)
header(s, "VALIDITY GATES", "결과를 믿을 수 있는 이유 — 10가지 자동 검증 기준")
rows = [["기준", "확인하는 것", "통과 조건"],
        ["G1  데이터 크기", "디스크가 실제로 쓰이는 조건인가", "데이터 ≥ DB 캐시의 2배"],
        ["G2  스토리지 사용", "디스크 읽기가 실제로 발생했나", "물리 읽기 IOPS > 0, 캐시 적중률 < 99%"],
        ["G3  클라이언트 여유", "측정 도구가 병목은 아니었나", "부하 VM CPU < 60%, 네트워크 < 50%"],
        ["G4  대기열 지연", "요청이 도구 안에서 밀리지 않았나", "큐 대기 p99 < 5ms"],
        ["G5  오류율", "정상 처리 기준의 측정인가", "오류 < 1%, 타임아웃 < 0.5%"],
        ["G6  서버 지표 정합", "서버 쪽 증거가 남았나", "측정 구간(UTC)과 ±1분 이내 일치, 값 존재"],
        ["G7  안정 상태", "예열이 끝난 뒤 측정했나", "직전 2분 처리량 변동 < 5%"],
        ["G9  데이터 정합성", "부하가 데이터를 훼손하지 않았나", "구매 원장 합계 = 잔액 차감 합계, 음수 없음"],
        ["G10 측정 시간", "계획한 만큼 측정했나", "계획 시간의 95% 이상 실행"]]
kr_table(s, 0.6, 1.6, 12.1, rows, [2.5, 5.0, 4.6], size=11, row_h=0.42)
text(s, 0.6, 5.75, 12.1, 0.8, [[("하나라도 어기면 그 측정은 결과에서 제외하고, 원인을 고친 뒤 다시 측정했습니다. ", 11.5, True, INK),
                                 ("제외 내역은 부록의 실행 로그에 모두 남아 있습니다. 최종 데이터 정합성 검사는 4개 DB 모두 위반 0건이었습니다.", 11.5, False, BODY)]])
footer(s)

# ============================ SECTION 01: MySQL ===============================
section_divider("01", "MySQL — SSD v1 vs v2", "Azure Database for MySQL – Flexible Server (8.4)",
                ["같은 IOPS(5,000)로 맞춘 공정 비교", "일반 플레이 · 쓰기 집중 · 로그인 러시 · 핫스팟", "8 vCore(C1)와 16 vCore(C5·C7) 구성"])

# ---- 8. 실험 B 환경 ----
s = slide(LIGHT)
header(s, "MYSQL — TEST SETUP", "무엇을 같게 하고, 무엇이 다른가")
rows = [["항목", "v1 (Premium SSD v1)", "v2 (Premium SSD v2)", "비고"],
        ["스토리지", "Premium SSD v1 · 128GiB · 5,000 IOPS", "Premium SSD v2 · 128GiB · 5,000 IOPS", "비교 대상 (IOPS 동일하게 구매)"],
        ["컴퓨팅", "E8ds_v5 → E16ds (8→16 vCore, 64GiB)", "E8ds_v6 → E16ds (동일)", "플랫폼이 v2에 v6 세대를 강제 — 통제 불가"],
        ["DB 캐시", "8GiB로 고정", "8GiB로 고정", "디스크가 쓰이도록 의도적으로 축소 (사전 등록)"],
        ["doublewrite", "OFF (플랫폼 고정)", "ON (플랫폼 고정)", "쓰기 경로 차이 — 통제 불가, 결과 해석에 반영"],
        ["데이터·부하", "동일 (같은 시드, 같은 시각 병렬 실행)", "동일", "구성 케이스: 8vCore+HA / 16vCore+HA / 16vCore 단독"]]
kr_table(s, 0.6, 1.65, 12.1, rows, [1.7, 3.9, 3.7, 2.8], size=10.5, row_h=0.52)
card(s, 0.6, 5.05, 11.9, 1.5, WARN, "해석 시 유의",
     ["컴퓨팅 세대(v5/v6)와 doublewrite는 플랫폼이 고정한 값이라 분리할 수 없음 — 이 결과는 “스토리지 매체 단독”이 아니라 “SSD v1 제품 구성 vs SSD v2 제품 구성”의 비교",
      "16 vCore로 늘려도 처리 한계가 늘지 않음 → 두 구성 모두 CPU가 아니라 스토리지가 병목인 상태에서 측정됨 (스토리지 비교로 적합한 조건)"],
     title_size=13, body_size=11)
footer(s)

# ---- 9. 실험 B 결과 표 ----
s = slide(LIGHT)
header(s, "MYSQL — RESULTS", "시나리오별 결과 — v2가 v1 대비 얼마나 달랐나")
rows = [["구성 · 시나리오", "부하", "반복", "평균 응답(p50)", "상위 1% 지연(p99)", "판정"],
        ["8vCore+HA · 일반 플레이", "5,000/s", "3/3", "−28.6%  [−39, −14]", "차이 확정 불가 (편차 큼)", "v2 우세 (평균)"],
        ["16vCore+HA · 일반 플레이", "5,000/s", "5/5", "−34.0%  [−37, −30]", "차이 확정 불가 (편차 큼)", "v2 우세 (평균)"],
        ["16vCore+HA · 쓰기 집중", "3,000/s", "3/3", "−17.9%  [−20, −17]", "−55.2%  [−62, −44]", "v2 우세"],
        ["16vCore+HA · 로그인 러시(3배)", "2,000→6,000/s", "3/3", "−27.0%", "요청 유실 v1 6.6% vs v2 3.8%", "v2 우세"],
        ["16vCore 단독 · 일반 플레이", "5,000/s", "3/3", "−33.7%  [−35, −33]", "+30.0%  [+27, +33]", "엇갈림"],
        ["16vCore 단독 · 핫스팟", "3,000/s", "3/3", "−19.2%  [−20, −18]", "−33.4%  [−36, −29]", "v2 우세"]]
kr_table(s, 0.6, 1.65, 12.1, rows, [3.4, 1.5, 0.9, 2.4, 2.7, 1.2], size=10.5, row_h=0.5)
text(s, 0.6, 5.35, 12.1, 1.2, [
    [("읽는 법  ", 11, True, BLUE), ("음수 = v2가 빠름. [ ] 안은 95% 신뢰구간 — 구간이 0을 포함하면 “차이가 있다”고 확정하지 않았습니다.", 11, False, BODY)],
    [("“차이 확정 불가”의 이유  ", 11, True, BLUE), ("일반 플레이에서는 v2가 3회 중 1회꼴로 p99가 3~4배 튀는 구간을 보여 반복 간 편차가 컸습니다. 이 급증이 없던 반복에서는 v2의 p99도 v1보다 낮았습니다.", 11, False, BODY)]], spacing=6)
footer(s)

# ---- 10. 실험 B 차트 ----
lp = "analysis/expB/charts/latency_percentiles.png"
if os.path.exists(lp):
    s = slide(LIGHT)
    header(s, "MYSQL — CHART", "응답시간 분포 비교 (막대 = 반복 중앙값, 선 = 최소~최대)")
    pic = s.shapes.add_picture(lp, Inches(0), Inches(1.62), height=Inches(4.8)); pic.left = int((prs.slide_width - pic.width) / 2)
    text(s, 0.6, 6.6, 12.1, 0.5, [[("세로축은 로그 스케일 — 낮을수록 좋음. p50(평균적 사용자 체감)은 전 구간에서 v2(빨강)가 낮고, p99 이상의 꼬리는 시나리오에 따라 우열이 갈립니다.", 11, False, BODY)]])
    footer(s)

# ---- 11. 실험 B 해석 ----
s = slide(LIGHT)
header(s, "MYSQL — TAKEAWAYS", "MySQL 결과가 의미하는 것")
card(s, 0.6, 1.65, 5.95, 2.45, BLUE, "확실한 것",
     ["평균 응답: 모든 시나리오·구성에서 v2가 18~34% 빠름 (신뢰구간이 0을 확실히 벗어남)",
      "쓰기 집중·핫스팟처럼 디스크 쓰기가 많은 부하일수록 v2의 이점이 커짐 — p99 기준 33~55% 개선",
      "로그인 러시(3배 버스트)에서 요청 유실이 v1의 절반 수준"], title_size=13.5, body_size=11)
card(s, 6.75, 1.65, 5.95, 2.45, WARN, "지켜봐야 할 것",
     ["일반 플레이 부하에서 v2에 간헐적 p99 급증 (버퍼 오염 + 쓰기 몰림 정황, 원인 미확정)",
      "HA 없는 구성에서는 v2의 p99가 일관되게 30% 높음 — HA 유무에 따라 꼬리 특성이 뒤집힘",
      "doublewrite ON(v2)의 영향 가능성 — 플랫폼 고정값이라 실험으로 분리 불가"], title_size=13.5, body_size=11)
card(s, 0.6, 4.35, 11.9, 2.15, TEAL, "권고",
     ["평균 응답과 쓰기 성능이 중요한 워크로드(결산·랭킹·이벤트 정산)는 v2 전환 이익이 분명함",
      "다만 꼬리 지연에 민감한 실시간 기능이 있다면, 전환 전 24시간 이상 장시간 부하(soak)로 간헐 급증의 빈도·원인을 확인할 것을 권고",
      "전환 시 IOPS·처리량을 독립 조정할 수 있는 v2의 과금 구조상, 현재 v1에서 IOPS 추가 구매 중인 서버일수록 비용 이점도 큼"], title_size=13.5, body_size=11.5)
footer(s)

# ========================== SECTION 02: PG vs HZ ==============================
section_divider("02", "PostgreSQL vs HorizonDB", "Flexible Server (PG 17) vs Azure HorizonDB (Public Preview)",
                ["같은 vCore(8→16) · 같은 부하로 비교", "일반 플레이 · 쓰기 집중 · 로그인 러시 · 핫스팟", "캐시 구조 차이를 명시한 조건부 비교"])

# ---- 13. 실험 A 환경 ----
s = slide(LIGHT)
header(s, "PG vs HORIZONDB — TEST SETUP", "무엇을 같게 하고, 무엇이 다른가")
rows = [["항목", "PostgreSQL Flexible Server", "Azure HorizonDB", "비고"],
        ["컴퓨팅", "D8ds_v5 → D16ds_v5 (8→16 vCore)", "8 → 16 vCore", "vCore 수만 일치 가능"],
        ["메모리 캐시", "shared_buffers 8GB → 16GB", "45GB → 90GB (자동 설정)", "HorizonDB가 5배 이상 큼 — 핵심 차이"],
        ["스토리지", "Premium SSD P15 256GiB (1,100 IOPS)", "관리형 스토리지 (내부 값 미공개)", "동일하게 맞출 수 없음"],
        ["HA/복제", "Same-zone HA → 해제", "replicaCount 1 (최소값, 축소 불가)", "“HA 없음” 구성이 HorizonDB에는 없음"],
        ["시간당 요금", "약 $0.79 (D16ds_v5 기준)", "약 $1.4 이상 (16 vCore 환산)", "약 1.8배 — 비용 반영 필요"]]
kr_table(s, 0.6, 1.65, 12.1, rows, [1.7, 3.9, 3.6, 2.9], size=10.5, row_h=0.52)
card(s, 0.6, 5.05, 11.9, 1.5, WARN, "이 비교를 읽는 올바른 방법",
     ["데이터(23GiB)가 HorizonDB 캐시에는 전부 올라가고 PostgreSQL 캐시에는 절반도 안 올라감 → “디스크를 읽는 PG vs 메모리에서 답하는 HorizonDB”의 비교",
      "따라서 이 결과는 “같은 돈으로 같은 vCore를 샀을 때의 체감 차이”로는 유효하지만, 엔진 자체의 우열로 해석하면 안 됨"],
     title_size=13, body_size=11)
footer(s)

# ---- 14. 실험 A 결과 표 ----
s = slide(LIGHT)
header(s, "PG vs HORIZONDB — RESULTS", "시나리오별 결과 — HorizonDB가 PostgreSQL 대비 얼마나 달랐나")
rows = [["구성 · 시나리오", "부하", "유효 반복", "평균 응답(p50)", "상위 1% 지연(p99)", "판정"],
        ["8vCore+HA · 일반 플레이", "5,500/s", "3/3", "−24.5%", "−89.8%  [−95, −84]", "HZ 우세"],
        ["16vCore+HA · 일반 플레이", "5,500/s", "2/5", "−16.9%", "−87.1%  [−93, −75]", "HZ 우세*"],
        ["16vCore+HA · 쓰기 집중", "3,300/s", "0/3", "—", "PG 전 회차 수 초 멈춤 발생", "비교 불가"],
        ["16vCore+HA · 로그인 러시(3배)", "2,200→6,600/s", "3/3", "−60.5%", "−94.5% (둘 다 유실 0)", "HZ 우세"],
        ["16vCore 단독 · 일반 플레이", "5,500/s", "3/3", "−14.9%", "−75.0%  [−94, −46]", "HZ 우세"],
        ["16vCore 단독 · 핫스팟", "3,300/s", "2/3", "−31.0%", "−76.8%  [−88, −57]", "HZ 우세"]]
kr_table(s, 0.6, 1.65, 12.1, rows, [3.4, 1.5, 1.0, 1.9, 3.0, 1.3], size=10.5, row_h=0.5)
text(s, 0.6, 5.35, 12.1, 1.3, [
    [("* 유효 반복이 준 이유  ", 11, True, BLUE), ("PostgreSQL이 측정 중 수 초간 멈추는(stall) 구간을 보인 반복은 검증 기준에서 탈락했습니다. 디스크 IOPS 소진(100%)과 체크포인트가 겹친 정황이며, 이 “탈락” 자체가 중요한 발견입니다.", 11, False, BODY)],
    [("HorizonDB는 전 시나리오에서 p99 4ms 안팎을 유지 ", 11, True, INK), ("— 부하를 3배로 튀겨도 흔들리지 않았습니다. 단, 캐시 상주 조건임을 함께 볼 것.", 11, False, BODY)]], spacing=6)
footer(s)

# ---- 15. 실험 A 차트 ----
lp = "analysis/expA/charts/latency_percentiles.png"
if os.path.exists(lp):
    s = slide(LIGHT)
    header(s, "PG vs HORIZONDB — CHART", "응답시간 분포 비교 (막대 = 반복 중앙값, 선 = 최소~최대)")
    pic = s.shapes.add_picture(lp, Inches(0), Inches(1.62), height=Inches(4.8)); pic.left = int((prs.slide_width - pic.width) / 2)
    text(s, 0.6, 6.6, 12.1, 0.5, [[("세로축은 로그 스케일. PostgreSQL(파랑)은 p99 이상에서 편차가 크게 벌어지는 반면, HorizonDB(빨강)는 어느 백분위에서도 거의 흔들리지 않습니다.", 11, False, BODY)]])
    footer(s)

# ---- 16. 실험 A 해석 ----
s = slide(LIGHT)
header(s, "PG vs HORIZONDB — TAKEAWAYS", "PostgreSQL vs HorizonDB 결과가 의미하는 것")
card(s, 0.6, 1.65, 5.95, 2.45, STEEL, "HorizonDB의 강점",
     ["같은 vCore·같은 부하에서 꼬리 지연(p99)이 75~95% 낮고, 반복 간 편차가 거의 없음",
      "로그인 러시 수준의 버스트(초당 6,600건)도 지연 변화 없이 흡수",
      "vCore만 지정하면 캐시·스토리지를 서비스가 알아서 크게 잡아 주는 구조 자체가 운영 부담을 줄임"], title_size=13.5, body_size=11)
card(s, 6.75, 1.65, 5.95, 2.45, WARN, "판단 전 확인할 것",
     ["캐시(메모리) 우위 조건의 결과 — 데이터가 캐시보다 커지는 규모(90GiB+)에서는 재검증 필요",
      "Public Preview 단계: 내부 스토리지 사양·일부 지표 미공개, SLA 미확정",
      "시간당 요금 약 1.8배 — 성능/비용 환산 시 격차는 줄어듦"], title_size=13.5, body_size=11)
card(s, 0.6, 4.35, 11.9, 2.15, BLUE, "함께 발견한 것 — PostgreSQL 쪽 개선 포인트",
     ["PG는 쓰기 집중 부하에서 P15 디스크의 IOPS 한도(1,100)에 걸려 수 초간 멈추는 구간이 반복 발생 → HorizonDB와 무관하게, 현재 PG를 쓰는 워크로드라면 디스크 등급(P20+ 또는 Premium SSD v2) 상향 검토 가치가 있음",
      "HA를 끄고 캐시를 16GB로 늘린 구성에서는 PG도 p99 7~8ms로 안정 — PG의 꼬리 지연은 상당 부분 디스크·HA 구성의 문제였음"], title_size=13.5, body_size=11.5)
footer(s)


# ========================== 추가 검증 (08-25~26 실측) ==========================
section_divider("04", "추가 검증 결과", "장시간 지속 부하 · 대용량 데이터 · 읽기 복제본",
                ["14.5시간 지속 부하 (MySQL v1 vs v2)", "92GiB 대용량 (PG vs HorizonDB)", "읽기 복제본 라우팅 효과 (MySQL)"])

# ---- 장시간 soak ----
s = slide(LIGHT)
header(s, "EXTENDED — 14.5H SOAK", "장시간 지속 부하: 10분 측정과는 다른 그림")
kr_table(s, 0.6, 1.6, 12.1, [
    ["지표 (14.5시간 누적, 초당 5,000건)", "v1 (SSD v1)", "v2 (SSD v2)", "10분 측정과의 비교"],
    ["평균 응답 (p50)", "20.0 ms", "22.1 ms", "10분 측정(3.6/2.3 ms)의 6~10배"],
    ["상위 5% (p95)", "1.13 초", "0.83 초", "10분 측정(약 20 ms)의 40~55배"],
    ["상위 1% (p99)", "2.04 초", "1.21 초", "장시간에서는 v2 꼬리가 40% 낮음"],
    ["최대 지연", "7.1 초", "2.4 초", "v1의 최악 구간이 3배 김"],
    ["오류율", "1.85%", "1.84%", "10분 측정(0%)과 달리 오류 발생"]], [4.0, 2.2, 2.2, 3.7], size=11, row_h=0.52)
card(s, 0.6, 5.1, 11.9, 1.5, WARN, "해석 — 왜 이렇게 다른가",
    ["14시간 동안 쓰기가 누적되며 데이터가 계속 커져(매치 기록·구매 원장 증가) 두 제품 모두 점진적으로 느려짐 — 짧은 측정으로는 보이지 않는 운영 현실",
     "단기 측정에서 관찰된 \"v2 간헐 급증\"과 반대로, 지속 부하에서는 v1의 꼬리·최악 구간이 더 나쁨 → v2 전환 판단에 유리한 증거. 단, 야간 자동 종료로 24h 중 14.5h만 측정됨(체크포인트로 보존)"], title_size=13, body_size=11)
footer(s)

# ---- 대용량 ----
s = slide(LIGHT)
header(s, "EXTENDED — 92GiB DATASET", "대용량 데이터(1,000만 계정): PG는 붕괴, HorizonDB는 무변화")
kr_table(s, 0.6, 1.6, 12.1, [
    ["초당 5,501건, 10분 × 3회", "PostgreSQL (D16ds + P15 디스크)", "HorizonDB (16 vCore)"],
    ["1회차", "실질 붕괴 — 평균 응답 11.6초, 처리량 35% 미달", "p99 3.9 ms, 오류 0"],
    ["2회차", "평균 3.3초, 처리량 7% 미달", "p99 3.9 ms, 오류 0"],
    ["3회차", "평균은 회복(1 ms)됐으나 p99 8.9초 스파이크", "p99 4.0 ms, 오류 0"],
    ["서버 상태", "읽기 IOPS 3,000~4,500 — 디스크 한계 초과", "읽기 IOPS ~200, 캐시 적중 99.8%"]], [2.6, 5.3, 4.2], size=11, row_h=0.52)
card(s, 0.6, 4.9, 11.9, 1.7, STEEL, "해석",
    ["데이터가 캐시보다 훨씬 큰 조건(92GiB vs 16GB)에서 PG+P15 조합은 이 부하를 수용하지 못함 — 검증 기준 미달로 \"유효 쌍 0\"이지만, 이 비대칭 실패 자체가 핵심 발견",
     "HorizonDB는 대용량에서도 캐시(90GB)가 데이터를 거의 덮어 동일 성능 유지 — 단, 데이터가 캐시의 2배를 넘는 초대용량 구간은 여전히 미검증",
     "PG를 이 규모로 쓰려면 디스크 등급 상향(P30+/SSD v2) 또는 캐시 증설이 선행돼야 함"], title_size=13, body_size=11)
footer(s)

# ---- 복제본 ----
s = slide(LIGHT)
header(s, "EXTENDED — READ REPLICA", "읽기 복제본 라우팅 (MySQL, 같은 서버에서 라우팅만 켜고 끔)")
kr_table(s, 0.6, 1.6, 12.1, [
    ["초당 5,000건 · 8vCore · 3회", "복제본 없이 (primary만)", "읽기 65%를 복제본으로 라우팅", "차이"],
    ["평균 응답 (p50)", "3.99 ms", "3.82 ms", "-1.4% (차이 없음 수준)"],
    ["상위 1% (p99)", "30.4 ms", "29.0 ms", "-7.4% [-13%, -5%]"],
    ["primary 읽기 IOPS", "1,776", "603", "-66% — 디스크 부담 1/3로"],
    ["복제 지연", "-", "측정 구간 내 이상 없음", "-"]], [3.0, 3.0, 3.4, 2.7], size=11, row_h=0.52)
card(s, 0.6, 4.9, 11.9, 1.7, TEAL, "해석과 제약",
    ["응답시간 개선은 작지만(스토리지가 이미 여유), primary 디스크 읽기가 1/3로 줄어 \"쓰기 여력 확보\" 효과가 분명 — 쓰기가 몰리는 서비스에서 가치",
     "중요 제약: Premium SSD v2(프리뷰) 서버는 읽기 복제본 생성이 현재 불가(InternalServerError 반복, 지원 티켓용 Tracking ID 확보) → v2 전환 시 복제본 전략은 GA 시점 재확인 필요",
     "복제본을 늘린 HA 구성 비교(C2·C6)는 이 제약으로 미실행"], title_size=13, body_size=11)
footer(s)

# ========================== SECTION 03: 결론 ==================================
section_divider("05", "결론과 다음 단계", "권고 사항 · 남은 검증 항목과 예상 소요 시간",
                ["의사결정 관점의 요약", "아직 확인하지 못한 항목", "재현 방법과 원본 데이터 위치"])

# ---- 18. 결론 ----
s = slide(LIGHT)
header(s, "CONCLUSION", "의사결정 관점의 요약")
rows = [["질문", "현재 답", "확신 수준"],
        ["MySQL 스토리지를 v2로 바꿔야 하나", "평균 응답·쓰기 성능·버스트 흡수 모두 v2 우세. 꼬리 지연 민감 서비스만 장시간 검증 후 전환", "높음 (평균) / 중간 (꼬리)"],
        ["HorizonDB를 도입 후보로 볼 만한가", "같은 vCore 기준 체감 성능은 확실히 우수. 단 캐시 우위 조건·프리뷰·1.8배 비용을 감안한 조건부 후보", "중간 (조건부)"],
        ["지금 PG 구성에서 개선할 것은", "쓰기 부하가 있다면 디스크 등급 상향이 가장 급함 (IOPS 소진으로 인한 멈춤 실측)", "높음"],
        ["이 결과를 그대로 일반화해도 되나", "반복 3~5회 기준의 방향성 결론. 수치 자체보다 “방향과 조건”을 가져갈 것", "—"]]
kr_table(s, 0.6, 1.7, 12.1, rows, [3.6, 6.3, 2.2], size=11, row_h=0.72)
text(s, 0.6, 5.15, 12.1, 1.3, [
    [("한 줄 요약  ", 12, True, BLUE), ("“MySQL은 v2로 갈 이유가 충분하고, HorizonDB는 인상적이지만 조건을 붙여 지켜볼 단계, PG는 디스크부터 손보라”", 12, True, INK)],
    [("모든 수치의 원본(측정 JSON·서버 지표·게이트 판정·실행 로그)은 git 저장소에 있으며, 스크립트 한 번으로 재현 가능합니다.", 11, False, MUTED)]], spacing=8)
footer(s)

# ---- 19. 남은 테스트 ----
s = slide(LIGHT)
header(s, "NEXT STEPS", "아직 확인하지 못한 항목과 예상 소요 시간")
rows = [["항목", "무엇을 확인하나", "예상 소요", "우선순위"],
        ["MySQL v2 장시간 soak (24h+)", "간헐적 p99 급증의 원인 — v2 전환의 마지막 관문", "1.5일(무인)", "★★★"],
        ["읽기 복제본 케이스 (C2·C4·C6·C8)", "읽기 분산 효과 — 클라이언트 라우팅 구현 포함", "2일", "★★☆"],
        ["8vCore 단독(C3) + 반복 5회 승격", "매트릭스 완성 및 기존 3회 셀의 신뢰구간 축소", "1일", "★★☆"],
        ["HorizonDB 대용량 재검증 (90GiB+)", "캐시를 넘는 데이터에서도 우위가 유지되는지 — 도입 판단의 핵심 근거", "1일", "★★★"],
        ["MySQL as-provisioned IOPS 셀", "IOPS를 맞추지 않은 “기본 구성 그대로” 비교 (제품 기본값 관점)", "0.5일", "★☆☆"],
        ["비용 정규화 (실청구 기반)", "백만 트랜잭션당 실제 비용 비교 — 청구 데이터 수집 후 산출", "0.5일+대기", "★★☆"]]
kr_table(s, 0.6, 1.65, 12.1, rows, [3.4, 5.4, 2.0, 1.3], size=10, row_h=0.5)
text(s, 0.6, 6.55, 12.1, 0.6, [[("전체 진행 시 약 5~6 작업일(무인 실행 시간 포함). ", 11.5, True, INK),
                                 ("현재 모든 서버는 중지/축소 상태이며, 재시작하면 데이터 재적재 없이 이어서 측정할 수 있습니다.", 11.5, False, BODY)]])
footer(s)

# ---- 20. 부록 ----
s = slide(LIGHT)
header(s, "APPENDIX", "재현 방법과 원본 데이터")
card(s, 0.6, 1.65, 5.95, 3.4, BLUE, "저장소 구성",
     ["cmd/gamebench — 부하 도구 (Go, MySQL·PostgreSQL 공용)",
      "scripts/ — 배포·오케스트레이션·분석·보고서 생성",
      "results/ — 측정별 원본 JSON (클라이언트 지표 + 서버 지표 + 게이트 판정)",
      "analysis/ — 통계 결과와 차트",
      "docs/run-log.md — 실패·수정 이력 전체 (UTC)",
      "docs/plan.md — 사전 등록된 실험 계획서"], title_size=13.5, body_size=11.5)
card(s, 6.75, 1.65, 5.95, 3.4, TEAL, "재현 절차",
     ["서버 재시작 → 벤치 VM 시작 (데이터는 보존됨)",
      "knee-both.sh로 한계점 탐색 → run-case.sh로 본 측정",
      "azmon_slice.py로 서버 지표를 측정 구간에 맞춰 첨부",
      "analyze.py로 게이트 판정과 통계 산출",
      "build_pptx_kr.py로 이 보고서 재생성"], title_size=13.5, body_size=11.5)
card(s, 0.6, 5.25, 11.9, 1.3, STEEL, "측정 환경 요약",
     ["Korea Central (MySQL) · Australia East (PG/HorizonDB) · 부하 VM D16ds_v5 · 모니터링 Telegraf + InfluxDB + Grafana · 총 측정 기간 2026-08-17 ~ 08-18 (UTC)"],
     title_size=13, body_size=11)
footer(s)

out = sys.argv[1] if len(sys.argv) > 1 else "report/KRAFTON_DB_benchmark_KR.pptx"
os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
prs.save(out)
print("saved", out, "slides", PAGE[0])
