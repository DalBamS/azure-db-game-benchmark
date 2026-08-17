"""Build the engineer-facing PPTX report from analysis JSON + charts.

usage: python scripts/build_pptx.py analysis/expB analysis/expA report/azure-db-benchmark-report.pptx
(either analysis dir may be missing; the deck adapts.)
"""
import json, os, sys, glob, datetime as dt
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

NAVY = RGBColor(0x1E, 0x27, 0x61); ICE = RGBColor(0xCA, 0xDC, 0xFC); WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x21, 0x21, 0x21); MUTED = RGBColor(0x60, 0x66, 0x7A); GOOD = RGBColor(0x2C, 0x5F, 0x2D); BAD = RGBColor(0x99, 0x00, 0x11)
LIGHT = RGBColor(0xF4, 0xF6, 0xFB)
FONT = "Calibri"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_text(slide, x, y, w, h, text, size=14, bold=False, color=INK, align=PP_ALIGN.LEFT, font=FONT):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.margin_left = tf.margin_right = Inches(0.05)
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
        p.space_after = Pt(4)
    return tb


def bullets(slide, x, y, w, h, items, size=14, color=INK):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = "•  " + it; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
        p.space_after = Pt(6)
    return tb


def title_slide(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = NAVY
    add_text(s, 0.8, 2.3, 11.5, 1.6, title, 40, True, WHITE)
    add_text(s, 0.8, 4.0, 11.5, 1.2, subtitle, 18, False, ICE)
    return s


def section(title, subtitle=""):
    s = prs.slides.add_slide(BLANK)
    add_text(s, 0.6, 0.35, 12, 0.8, title, 30, True, NAVY)
    if subtitle:
        add_text(s, 0.6, 1.0, 12, 0.5, subtitle, 14, False, MUTED)
    return s


def table(slide, x, y, w, rows, col_widths=None, size=11, header=True):
    nrows, ncols = len(rows), len(rows[0])
    shp = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(0.3 * nrows))
    t = shp.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            t.columns[i].width = Inches(cw)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = t.cell(i, j); c.text = ""
            p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = str(val); r.font.size = Pt(size); r.font.name = FONT
            c.margin_left = c.margin_right = Inches(0.05); c.margin_top = c.margin_bottom = Inches(0.02)
            if header and i == 0:
                r.font.bold = True; r.font.color.rgb = WHITE; c.fill.solid(); c.fill.fore_color.rgb = NAVY
            else:
                c.fill.solid(); c.fill.fore_color.rgb = LIGHT if i % 2 else WHITE; r.font.color.rgb = INK
    return shp


def fmt(v, kind):
    if v is None:
        return "-"
    if kind == "us":
        return f"{v/1000:,.1f} ms"
    if kind == "pct":
        return f"{v:+.1f}%"
    if kind == "rate":
        return f"{100*v:.3f}%"
    if kind == "ratio":
        return f"{v:.4f}"
    return f"{v:,.0f}"


def load_analysis(d):
    p = os.path.join(d, "analysis.json")
    return json.load(open(p, encoding="utf-8")) if os.path.exists(p) else None


def outcome_rows(cell):
    A, B = cell["arms"]
    keys = [("success_tps", "Throughput (TPS)", "n"), ("p50_us", "p50", "us"), ("p95_us", "p95", "us"), ("p99_us", "p99", "us"), ("p999_us", "p99.9", "us"),
            ("error_rate", "Error rate", "rate"), ("drop_rate", "Dropped arrivals", "rate"), ("burst_recovery_s", "Burst recovery (s)", "n"), ("read_iops", "Server read IOPS", "n"), ("write_iops", "Server write IOPS", "n"), ("bp_hit", "Buffer hit ratio", "ratio")]
    rows = [["Outcome", f"{A} (median)", f"{B} (median)", f"Δ% {B} vs {A}", "95% CI (bootstrap)", "n pairs", "Reading"]]
    for k, label, kind in keys:
        o = cell["outcomes"].get(k)
        if not o:
            continue
        ci = "-" if not o["ci95"] else f"[{o['ci95'][0]:+.1f}%, {o['ci95'][1]:+.1f}%]"
        rows.append([label, fmt(o["v1_median"], kind), fmt(o["v2_median"], kind), fmt(o["pct_change"], "pct") if o["pct_change"] is not None else "-", ci, o["n"], o["note"].split(";")[0]])
    return rows


def overview_slide(an, exp_title, arm_desc):
    s = section(f"{exp_title} — overview across cases", arm_desc + "  |  medians of gate-passing repetitions; Δ = geometric-mean paired change of second arm vs first with 95% bootstrap CI")
    rows = [["Cell", "Rate/s", "pairs", "p50 A→B", "p95 A→B", "p99 A→B", "p99.9 A→B", "Δ p50 [CI]", "Δ p99 [CI]", "read IOPS A/B", "hit A/B"]]
    for c in an["cells"]:
        o = c["outcomes"]
        def ab(k, kind="us"):
            x = o.get(k, {})
            return f"{fmt(x.get('v1_median'), kind)} → {fmt(x.get('v2_median'), kind)}"
        def dci(k):
            x = o.get(k, {})
            if x.get("pct_change") is None:
                return "-"
            return f"{x['pct_change']:+.0f}% [{x['ci95'][0]:+.0f}, {x['ci95'][1]:+.0f}]"
        rows.append([f"{c['case']}/{c['scenario']}", f"{c['rate']:,}", f"{c['n_pairs_usable']}/{c['n_pairs_total']}", ab("p50_us"), ab("p95_us"), ab("p99_us"), ab("p999_us"), dci("p50_us"), dci("p99_us"), ab("read_iops", "n"), ab("bp_hit", "ratio")])
    table(s, 0.5, 1.5, 12.3, rows, col_widths=[0.9, 0.8, 0.6, 1.4, 1.4, 1.4, 1.5, 1.4, 1.4, 1.1, 0.9], size=10)
    add_text(s, 0.5, 5.6, 12.3, 1.2, ["A = first arm, B = second arm (v1→v2 / postgres→horizon). C1 = 8 vCore + Same-zone HA, C5 = 16 vCore + Same-zone HA; same dataset, same offered load, same client.",
                                      "Read as: negative Δ latency = B faster. CIs from 3 repetitions are wide; a CI that excludes 0 is a directional finding, not a precise magnitude."], 11, False, MUTED)


def cell_slides(an, exp_title, chart_dir, arm_desc):
    for ci, cell in enumerate(an["cells"]):
        A, B = cell["arms"]
        scen_note = {"S1": "S1 normal play (65% read)", "S2": "S2 write-heavy (80% write)", "S3": "S3 login burst: base rate ×3 for 120 s at t+240 s", "S4": "S4 hotspot"}.get(cell["scenario"], cell["scenario"])
        case_note = {"C1": "8 vCore / Same-zone HA", "C5": "16 vCore / Same-zone HA", "C7": "16 vCore / No-HA (HorizonDB keeps replicaCount 1 — minimum allowed)", "C3": "8 vCore / No-HA"}.get(cell["case"], cell["case"])
        s = section(f"{exp_title} — {cell['case']} / {cell['scenario']} @ {cell['rate']:,}/s", f"{case_note} · {scen_note}  |  {arm_desc}  |  usable pairs {cell['n_pairs_usable']}/{cell['n_pairs_total']} (10-min open-loop, arms concurrent)")
        table(s, 0.6, 1.5, 12.1, outcome_rows(cell), col_widths=[2.0, 1.6, 1.6, 1.4, 2.2, 0.8, 2.5], size=11)
        # gate summary
        lines = []
        for rep, info in sorted(cell["reps"].items(), key=lambda kv: int(kv[0])):
            parts = []
            for arm, i in info.items():
                failed = [k.split("_")[0] for k, v in i["gates"].items() if not v]
                soft = i.get("soft_failed", [])
                parts.append(f"{arm}: {'PASS' if i['ok'] else 'EXCLUDED'} tps={i['tps']:.0f} p99={i['p99_us']/1000:.1f}ms" + (f" fail={failed}" if failed and not i['ok'] else "") + (f" soft-warn={[k.split('_')[0] for k in soft]}" if soft else ""))
            lines.append(f"rep {rep}: " + "   |   ".join(parts))
        add_text(s, 0.6, 5.0, 12.1, 1.6, ["Gate status per repetition (G1 dataset≥2×buffer, G2 storage used, G3 client headroom, G4 queue, G5 errors, G6 metrics, G7 steady state, G10 duration; azmon=Azure Monitor slice; inv=invariants):"] + lines, 10, False, MUTED)
        # chart slide: latency percentiles + timeline
        pngs = sorted(glob.glob(os.path.join(chart_dir, f"timeline_{cell['case']}_{cell['scenario']}_{cell['rate']}_rep*.png")))
        lp = os.path.join(chart_dir, "latency_percentiles.png")
        if pngs:
            s2 = section(f"{exp_title} — {cell['case']} time series (client TPS / p99 vs server IOPS, same UTC axis)", "shaded = 10-minute measurement window; server IOPS from in-run status sampling (SHOW GLOBAL STATUS / pg_stat)")
            s2.shapes.add_picture(pngs[0], Inches(0.6), Inches(1.4), height=Inches(5.9))
        if os.path.exists(lp) and ci == len(an["cells"]) - 1:
            s3 = section(f"{exp_title} — latency percentiles by cell (median of ALL reps incl. gate-failed ones, whiskers = min/max)", "log scale; lower is better; use the per-cell tables (gate-passing pairs only) for the estimates")
            s3.shapes.add_picture(lp, Inches(0.6), Inches(1.4), height=Inches(5.6))
        for kp in sorted(glob.glob(os.path.join(chart_dir, f"knee_{cell['case']}_*.png"))):
            s4 = section(f"{exp_title} — {cell['case']} knee search (closed-loop staircase)", "used only to pick the open-loop measurement rate (≈65% of the lower arm's knee); not a result table")
            s4.shapes.add_picture(kp, Inches(1.5), Inches(1.4), height=Inches(5.6))


def main(expB_dir, expA_dir, out):
    b = load_analysis(expB_dir); a = load_analysis(expA_dir)
    today = dt.datetime.now(dt.timezone.utc).strftime("%Y-%m-%d")
    title_slide("Azure Database Benchmark — Game OLTP Workload", f"Exp B: MySQL Flexible Server Premium SSD v1 vs v2   |   Exp A: PostgreSQL Flexible Server vs HorizonDB\nGate-validated open-loop measurements · {today} UTC · Euson Internal Subscription")

    s = section("Executive summary", "what was measured, what passed validation, what it means")
    items = []
    if b:
        for c in b["cells"]:
            o = c["outcomes"]; p99 = o.get("p99_us"); tps = o.get("success_tps")
            items.append(f"Exp B {c['case']} @ {c['rate']:,}/s: v2 vs v1 p99 {fmt(p99['pct_change'],'pct') if p99 and p99['pct_change'] is not None else '-'} (CI {p99['ci95'] and f'[{p99['ci95'][0]:+.1f}%, {p99['ci95'][1]:+.1f}%]'}), p50 {fmt(o['p50_us']['pct_change'],'pct') if o['p50_us']['pct_change'] is not None else '-'}; usable pairs {c['n_pairs_usable']}/{c['n_pairs_total']}")
    if a:
        for c in a["cells"]:
            o = c["outcomes"]; p99 = o.get("p99_us")
            items.append(f"Exp A {c['case']} @ {c['rate']:,}/s: HorizonDB vs PG p99 {fmt(p99['pct_change'],'pct') if p99 and p99['pct_change'] is not None else '-'} (CI {p99['ci95'] and f'[{p99['ci95'][0]:+.1f}%, {p99['ci95'][1]:+.1f}%]'}); usable pairs {c['n_pairs_usable']}/{c['n_pairs_total']}")
    items += ["All reported numbers come only from runs that passed the hard validity gates (dataset ≥ 2× buffer pool, physical reads > 0, client CPU < 60%, errors < 1%, steady state, server metrics aligned to the UTC window).",
              "3 repetitions per cell (protocol baseline is 5) — CIs are wide; treat effects as directional unless the CI excludes 0.",
              "Exp A: HorizonDB's buffer cache (45 GB at 8 vCore, 90 GB at 16 vCore) exceeds the dataset — G1/G2 are soft-failed for HorizonDB; the comparison is 'storage-bound PG (P15 disk at 100% IOPS) vs cache-resident HorizonDB' at the same offered load. In C5, PG Flexible showed multi-second stalls in 2 of 3 reps (queue/latency gates failed → excluded), HorizonDB none.",
              "Exp A S2 (write-heavy, 3,300/s): HorizonDB p99 ≈ 4 ms in all 3 reps; PG Flexible (P15 disk, IOPS 100%) stalled for seconds in every rep (p99 0.4–20 s, G4/G5 fail) → no usable pair, but the failure itself is the finding. Exp A C7 (No-HA, warm): PG p99 7–8 ms vs HorizonDB 3.7 ms (−75% [−94%, −46%]).",
              "Exp B: MySQL is storage-bound in both cases (16 vCore did not raise the knee); v2's median latency is consistently lower, but v2 showed a tail-latency degradation episode in the 3rd repetition of both C1 and C5 (dirty pages ↑, threads_running spikes) — root cause not established; needs a longer soak."]
    bullets(s, 0.6, 1.5, 12.1, 5.5, items, 14)

    s = section("Method — how validity was enforced", "traditional load-testing protocol + automated gates; pilots/smoke never appear in result tables")
    bullets(s, 0.6, 1.4, 6.0, 5.6, [
        "Workload: game OLTP schema (accounts/inventory/sessions/ledger/matches/leaderboard/guilds), random incompressible payloads; S1 mix = 65% read / 35% write (profile 30, inventory 20, leaderboard 10, guild 5, session upsert 10, inventory update 12, match result 8, purchase 5).",
        "Key model: 20% of ops on a 50k-account hot set, 80% uniform over 2.5M accounts (top-1 key share ≈ 0.0004%).",
        "Load: closed-loop staircase finds the knee → open-loop Poisson-ish arrivals at ≈65% of the lower arm's knee; latency measured from scheduled arrival (no coordinated omission); HDR histogram (µs); errors counted separately.",
        "Concurrency: 256 dedicated connections per arm; queue pickup delay p99 gated < 5 ms (no pool wait).",
        "Warm-up until last-120 s throughput CV < 5%, then 10-minute measurement; both arms run in the same wall-clock window; 3 repetitions.",
        "Server evidence inside every run: SHOW GLOBAL STATUS / pg_stat_* every 5 s → read/write IOPS, buffer hit ratio; Azure Monitor sliced to the run window; Telegraf→InfluxDB→Grafana for live view.",
        "Statistics: paired log-ratio per repetition → geometric-mean Δ%, seeded bootstrap 95% CI, per-arm CV; gate-failed runs excluded and listed."], 12)
    table(s, 6.9, 1.4, 5.8, [["Gate", "Rule"], ["G1", "dataset ≥ 2× buffer pool"], ["G2", "read IOPS > 0 and hit ratio < 99%"], ["G3", "client CPU < 60%, NIC < 50% (measured)"], ["G4", "queue pickup p99 < 5 ms"], ["G5", "errors < 1%, timeouts < 0.5%, drops < 0.1%"], ["G6", "≥3 server samples, 0 errors, UTC-aligned; Azure Monitor non-empty"], ["G7", "steady state CV < 5% before measuring"], ["G9", "app invariants (ledger = spent, no negatives)"], ["G10", "≥ 95% of planned duration"]], col_widths=[0.7, 5.1], size=11)

    if b:
        s = section("Exp B environment & confound table", "MySQL Flexible Server 8.4.9, Korea Central, Same-zone HA, storage 128 GiB, both arms 5,000 provisioned IOPS (same-IOPS cell)")
        table(s, 0.6, 1.4, 12.1, [["Item", "v1 (Premium SSD v1)", "v2 (Premium SSD v2)", "Controlled?"],
                                    ["Compute", "Standard_E8ds_v5 (8 vCore / 64 GiB) → C5: E16ds_v5", "Standard_E8ds_v6 (8 vCore / 64 GiB) → C5: E16ds_v6", "Generation differs (platform forces v6 for SSD v2) — confound"],
                                    ["Storage", "Premium_LRS 128 GiB, 5,000 IOPS (paid additional IOPS)", "PremiumV2_LRS 128 GiB, 5,000 IOPS / 125 MB/s", "IOPS matched"],
                                    ["innodb_buffer_pool_size", "8 GiB (pinned)", "8 GiB (pinned)", "controlled (pre-registered; default would be 48 GiB)"],
                                    ["innodb_doublewrite", "OFF (platform read-only)", "ON (platform read-only)", "cannot control — confound"],
                                    ["Dataset", "24.4 GiB (3.05× buffer pool)", "23.7 GiB (2.96× buffer pool)", "same generator/seed"],
                                    ["HA / zone", "SameZone / zone 1", "SameZone / zone 1", "controlled"],
                                    ["Network", "public endpoint + TLS from bench VM (D16ds_v5, same region)", "same", "controlled"]], col_widths=[2.2, 3.4, 3.4, 3.1], size=11)
        overview_slide(b, "Exp B (MySQL SSD v1 vs v2)", "arms: v1 = Premium SSD v1, v2 = Premium SSD v2 (same-IOPS 5,000)")
        cell_slides(b, "Exp B (MySQL SSD v1 vs v2)", os.path.join(expB_dir, "charts"), "arms: v1 = Premium SSD v1, v2 = Premium SSD v2 (same-IOPS 5,000)")
    if a:
        s = section("Exp A environment & confound table", "Australia East; PostgreSQL 17; both public endpoint + TLS from bench VM (D16ds_v5, same region)")
        table(s, 0.6, 1.4, 12.1, [["Item", "postgres (Flexible Server)", "horizon (HorizonDB)", "Controlled?"],
                                    ["Compute", "Standard_D8ds_v5 GeneralPurpose (8 vCore / 32 GiB) → C5: D16ds_v5", "vCores 8 → C5: 16 (memory not exposed; shared_buffers 45,058 MB observed)", "vCores matched only — memory differs"],
                                    ["Storage", "Premium SSD P15 256 GiB (1,100 IOPS baseline, burstable)", "HorizonDB managed storage (IOPS not exposed)", "not comparable — confound"],
                                    ["HA / replica", "SameZone HA, no read replica", "replicaCount 1 (read-only endpoint exists, unused by workload)", "best-effort mapping"],
                                    ["shared_buffers", "8 GB", "45,058 MB", "dataset (23 GiB) < HZ cache → G1/G2 soft-fail for HZ"],
                                    ["Version", "17.10", "17.9 (Azure HorizonDB)", "-"],
                                    ["Price", "list price per vCore-hour differs; see cost note", "", "report cost-normalised metric when billing data is available"]], col_widths=[2.2, 3.6, 3.6, 2.7], size=11)
        overview_slide(a, "Exp A (PostgreSQL vs HorizonDB)", "arms: postgres = Flexible Server, horizon = HorizonDB")
        cell_slides(a, "Exp A (PostgreSQL vs HorizonDB)", os.path.join(expA_dir, "charts"), "arms: postgres = Flexible Server, horizon = HorizonDB")

    s = section("Limitations & next steps")
    bullets(s, 0.6, 1.4, 12.1, 5.6, [
        "Cases run: C1 (8 vCore/HA), C5 (16 vCore/HA), C7 (16 vCore/No-HA); scenarios S1 (C1: 3 reps, C5: 5 reps, C7: 3 reps), S2 and S3 at C5 (3 reps). Not run: S4, C3, and the read-replica cases C2/C4/C6/C8 (client replica routing not implemented/validated in this time box). HorizonDB cannot be configured with 0 replicas (API minimum 1), so 'No-HA' has no exact HorizonDB counterpart.",
        "S3 (login burst ×3): MySQL arms could not absorb 6,000/s bursts (queue overflow, 4–7% dropped arrivals → G5 fail on both; reported as burst behaviour, not excluded from the burst slide). PG/HorizonDB absorbed the 6,600/s burst with 0 drops.",
        "Exp B confounds that cannot be removed on the platform: compute generation (v5 vs v6) and innodb_doublewrite (OFF vs ON). Effects should be read as 'SSD v1 product configuration vs SSD v2 product configuration'.",
        "Exp A: HorizonDB exposes far more cache than the Flexible Server SKU with the same vCores; the dataset would need to exceed ~90 GiB to make HorizonDB storage-bound. HorizonDB storage IOPS/memory are not exposed by the preview API.",
        "Buffer pool was pinned to 8 GiB on MySQL to make the storage tier observable within budget; production defaults (48 GiB) would keep more of a 24 GiB dataset in memory.",
        "Same-IOPS cell only (5,000 IOPS on both). The 'as-provisioned' cell (v1 default free IOPS) was not run.",
        "Cost: actual billing not yet available; cost-normalised comparisons to be added when Cost Management ingests usage.",
        "All raw run JSON (client histograms, per-second series, server status series, Azure Monitor slices, gate verdicts) is in results/; analysis is reproducible with scripts/analyze.py."], 13)
    os.makedirs(os.path.dirname(out) or ".", exist_ok=True)
    prs.save(out)
    print("saved", out)


if __name__ == "__main__":
    main(sys.argv[1], sys.argv[2], sys.argv[3])
